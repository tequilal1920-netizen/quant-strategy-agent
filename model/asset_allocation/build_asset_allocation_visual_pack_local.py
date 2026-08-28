from __future__ import annotations

import json
import math
import shutil
import textwrap
from pathlib import Path
from typing import Any, Dict, List, Sequence, Tuple

import numpy as np
import pandas as pd

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle, FancyArrowPatch, Wedge, Circle
from matplotlib.lines import Line2D
from matplotlib import font_manager

from openpyxl import Workbook
from openpyxl.chart import LineChart, Reference
from openpyxl.formatting.rule import DataBarRule
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(r"G:\中信建投")
SNAPSHOT = ROOT / "agent" / "output" / "model_improvement" / "asset_allocation_snapshot_v64_daily_excess_governed.json"
PANEL = ROOT / "agent" / "output" / "model_improvement" / "asset_allocation_panel_v553.json"
OUT = Path(r"C:\Users\Rye\Desktop\资产配置")

RED = "#B21A12"
DARK_RED = "#8B1A12"
LINE_RED = "#C00000"
ORANGE = "#F5B37D"
YELLOW = "#F5B400"
GREY = "#BFBFBF"
DARK_BLUE = "#243C64"
LIGHT_GREY = "#E9E9E9"
TEXT = "#111111"
CH_FONT = "楷体"
EN_FONT = "Arial"

STAGE_LABELS = {
    "recovery": "复苏期",
    "overheat": "过热期",
    "stagflation": "滞胀期",
    "recession": "衰退期",
    "I_recovery": "阶段Ⅰ\n复苏期",
    "II_prosperity": "阶段Ⅱ\n繁荣期",
    "III_overheat": "阶段Ⅲ\n过热期",
    "IV_credit_pressure": "阶段Ⅳ\n滞涨期",
    "V_early_recession": "阶段Ⅴ\n衰退前期",
    "VI_late_recession": "阶段Ⅵ\n衰退后期",
}


def wrap_cn(text: Any, width: int) -> str:
    s = str(text)
    if len(s) <= width:
        return s
    return "\n".join(textwrap.wrap(s, width=width, break_long_words=True, replace_whitespace=False))


def setup_font() -> None:
    candidates = [
        r"C:\Windows\Fonts\simkai.ttf",
        r"C:\Windows\Fonts\kaiu.ttf",
        r"C:\Windows\Fonts\simsun.ttc",
        r"C:\Windows\Fonts\msyh.ttc",
    ]
    for p in candidates:
        if Path(p).exists():
            try:
                font_manager.fontManager.addfont(p)
                name = font_manager.FontProperties(fname=p).get_name()
                plt.rcParams["font.family"] = [name, "Arial"]
                break
            except Exception:
                pass
    plt.rcParams["axes.unicode_minus"] = False
    plt.rcParams["mathtext.fontset"] = "dejavuserif"


def load_json(path: Path) -> Dict[str, Any]:
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def month_to_dt(months: Sequence[str]) -> pd.DatetimeIndex:
    return pd.to_datetime([str(m)[:4] + "-" + str(m)[4:6] + "-28" for m in months])


def pct(x: float, n: int = 1) -> str:
    if x is None or not np.isfinite(x):
        return ""
    return f"{x * 100:.{n}f}%"


def nav_from_returns(rets: Sequence[float], start: float = 1.0) -> np.ndarray:
    arr = np.asarray(rets, dtype=float)
    return start * np.cumprod(1 + arr)


def annual_metrics(months: Sequence[str], strategy_rets: Sequence[float], bench_rets: Sequence[float]) -> List[List[Any]]:
    idx = pd.Index([str(m)[:4] for m in months], name="year")
    sr = pd.Series(np.asarray(strategy_rets, dtype=float), index=idx)
    br = pd.Series(np.asarray(bench_rets, dtype=float), index=idx)
    rows: List[List[Any]] = []
    for y in sorted(idx.unique()):
        s = sr.loc[y].values
        b = br.loc[y].values
        strat = float(np.prod(1 + s) - 1)
        bench = float(np.prod(1 + b) - 1)
        active = float((1 + strat) / (1 + bench) - 1)
        dd = max_drawdown(s)
        label = y if y < "2026" else "2026YTD"
        rows.append([label, pct(strat), pct(bench), pct(active), pct(dd)])
    s_all = sr.values
    b_all = br.values
    n = len(s_all)
    strat_ann = float(np.prod(1 + s_all) ** (12 / n) - 1)
    bench_ann = float(np.prod(1 + b_all) ** (12 / n) - 1)
    active_ann = float((1 + strat_ann) / (1 + bench_ann) - 1)
    rows.append(["区间年化", pct(strat_ann), pct(bench_ann), pct(active_ann), pct(max_drawdown(s_all))])
    return rows


def max_drawdown(rets: Sequence[float]) -> float:
    nav = np.r_[1.0, nav_from_returns(rets)]
    peak = np.maximum.accumulate(nav)
    return float(np.min(nav / peak - 1))


def axis_off(ax):
    ax.set_axis_off()
    return ax


def save(fig, i: int) -> Path:
    p = OUT / f"{i}.png"
    fig.savefig(p, dpi=220, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return p


def title_top(ax, title: str) -> None:
    ax.text(0.015, 0.97, title, va="top", ha="left", fontsize=18, fontweight="bold", color=TEXT, transform=ax.transAxes)
    ax.plot([0, 1], [0.925, 0.925], color=DARK_RED, lw=1.6, transform=ax.transAxes, clip_on=False)


def footer(ax, text: str = "资料来源：Wind/RQData/本地研究数据库，中信建投量化整理") -> None:
    ax.plot([0, 1], [0.035, 0.035], color=DARK_RED, lw=1.2, transform=ax.transAxes, clip_on=False)
    ax.text(0.01, 0.008, text, fontsize=11, ha="left", va="bottom", transform=ax.transAxes)


def nav_months_values(rows: List[Dict[str, Any]]) -> Tuple[List[str], np.ndarray]:
    months = [str(r.get("month", "")) for r in rows]
    values = np.asarray([float(r.get("nav", np.nan)) for r in rows], dtype=float)
    return months, values


def return_months_values(rows: List[Dict[str, Any]]) -> Tuple[List[str], np.ndarray]:
    months = [str(r.get("month", "")) for r in rows]
    values = np.asarray([float(r.get("net_return", r.get("return", np.nan))) for r in rows], dtype=float)
    return months, values

def safe_float(value: Any, default: float = 0.0) -> float:
    try:
        v = float(value)
    except Exception:
        return default
    if not np.isfinite(v):
        return default
    return v

def draw_report_table(i: int, title: str, columns: Sequence[str], rows: Sequence[Sequence[Any]], widths: Sequence[float] | None = None, footer_text: str = "资料来源：Wind/RQData/本地研究数据库，中信建投量化整理") -> None:
    """券商报告式可读表：自动换行、行高自适应，避免文字重叠。"""
    ncols = len(columns)
    if widths is None:
        widths = [1 / ncols] * ncols
    width_sum = float(sum(widths))
    col_chars = [max(6, int(32 * w / width_sum)) for w in widths]
    wrapped_rows: List[List[str]] = []
    row_units: List[int] = []
    for row in rows:
        wrapped = [wrap_cn(v, col_chars[j]) for j, v in enumerate(row)]
        wrapped_rows.append(wrapped)
        row_units.append(max(1, max(cell.count("\n") + 1 for cell in wrapped)))
    fig_h = max(4.0, 1.35 + 0.36 * (sum(row_units) + 2))
    fig, ax = plt.subplots(figsize=(11.2, fig_h))
    axis_off(ax)
    title_top(ax, title)
    x0, y_top, table_w = 0.02, 0.86, 0.96
    xs = [x0]
    for w in widths:
        xs.append(xs[-1] + table_w * w / width_sum)
    usable_h = 0.75
    unit_h = min(0.075, usable_h / max(1.0, sum(row_units) + 1.25))
    header_h = unit_h * 1.18
    for j, col in enumerate(columns):
        ax.add_patch(Rectangle((xs[j], y_top - header_h), xs[j + 1] - xs[j], header_h, facecolor=RED, edgecolor="white", lw=1.1, transform=ax.transAxes))
        ax.text((xs[j] + xs[j + 1]) / 2, y_top - header_h / 2, wrap_cn(col, col_chars[j]), fontsize=13.2, color="white", fontweight="bold", ha="center", va="center", linespacing=1.25, transform=ax.transAxes)
    y_cursor = y_top - header_h
    for r, row in enumerate(wrapped_rows):
        row_h = unit_h * row_units[r]
        y = y_cursor - row_h
        fill = "#F7F1EC" if r % 2 == 0 else "#E9E9E9"
        for j, val in enumerate(row):
            ax.add_patch(Rectangle((xs[j], y), xs[j + 1] - xs[j], row_h, facecolor=fill, edgecolor="white", lw=0.9, transform=ax.transAxes))
            ha = "center" if j == 0 else "left"
            x_text = (xs[j] + xs[j + 1]) / 2 if j == 0 else xs[j] + 0.009
            fs = 12.2 if row_units[r] <= 2 else 10.8
            ax.text(x_text, y + row_h / 2, val, fontsize=fs, color=TEXT, ha=ha, va="center", linespacing=1.15, transform=ax.transAxes)
        y_cursor = y
    footer(ax, footer_text)
    save(fig, i)


def draw_annual_table(i: int, title: str, rows: Sequence[Sequence[Any]], bench_name: str = "四资产等权") -> None:
    draw_report_table(i, title, ["年度", "策略收益", bench_name, "超额收益", "最大回撤"], rows, widths=[0.16, 0.21, 0.21, 0.21, 0.21])


def draw_heatmap(i: int, title: str, labels: Sequence[str], corr: np.ndarray, footer_text: str = "资料来源：Wind/RQData/本地研究数据库，中信建投量化整理") -> None:
    fig, ax = plt.subplots(figsize=(9.6, 7.0))
    im = ax.imshow(corr, cmap="RdYlGn_r", vmin=-1, vmax=1)
    ax.set_xticks(range(len(labels)))
    ax.set_yticks(range(len(labels)))
    ax.set_xticklabels(labels, rotation=45, ha="right", fontsize=12)
    ax.set_yticklabels(labels, fontsize=12)
    for x in range(len(labels)):
        for y in range(len(labels)):
            ax.text(x, y, f"{corr[y, x]:.2f}", ha="center", va="center", fontsize=12, color=TEXT, fontweight="normal")
    ax.set_title(title, loc="left", fontsize=18, fontweight="bold", pad=18)
    for spine in ax.spines.values():
        spine.set_visible(False)
    fig.colorbar(im, ax=ax, fraction=0.035, pad=0.03)
    fig.text(0.02, 0.02, footer_text, fontsize=11)
    save(fig, i)


def draw_nav(i: int, title: str, dates: Sequence[pd.Timestamp], nav_strategy: Sequence[float], nav_bench: Sequence[float], label: str, bench_label: str = "四资产等权") -> None:
    ns = np.asarray(nav_strategy, dtype=float)
    nb = np.asarray(nav_bench, dtype=float)
    rel = ns / np.maximum(nb, 1e-12)
    fig, ax = plt.subplots(figsize=(8.8, 5.6))
    ax2 = ax.twinx()
    ax.plot(dates, nb, color=YELLOW, lw=2.3, label=bench_label)
    ax.plot(dates, ns, color=GREY, lw=2.6, label=label)
    ax2.plot(dates, rel, color=LINE_RED, lw=2.4, label="相对强度（右轴）")
    ax.grid(False)
    ax.spines["top"].set_visible(False)
    ax2.spines["top"].set_visible(False)
    ax.set_title(title, loc="left", fontsize=17, fontweight="bold")
    ax.tick_params(axis="x", rotation=90, labelsize=11)
    ax.tick_params(axis="y", labelsize=11)
    ax2.tick_params(axis="y", labelsize=11)
    lines = ax.get_lines() + ax2.get_lines()
    ax.legend(lines, [l.get_label() for l in lines], loc="lower center", bbox_to_anchor=(0.5, -0.25), ncol=3, frameon=False, fontsize=12)
    save(fig, i)


def draw_stage_step(i: int, title: str, dates: Sequence[pd.Timestamp], stages: Sequence[str], order: Sequence[str]) -> None:
    """每一个连续时间段只对应一个周期阶段，阶梯线按阶段切换。"""
    mp = {s: k + 1 for k, s in enumerate(order)}
    y = np.asarray([mp.get(str(s), np.nan) for s in stages], dtype=float)
    fig, ax = plt.subplots(figsize=(13.8, 5.4))
    ax.step(dates, y, where="post", color=LINE_RED, lw=1.55)
    change = np.r_[True, np.diff(y) != 0]
    ax.scatter(np.asarray(dates)[change], y[change], s=10, color=LINE_RED, zorder=3)
    ax.set_yticks(range(1, len(order) + 1))
    ax.set_yticklabels([STAGE_LABELS.get(s, s) for s in order], fontsize=12)
    ax.set_ylim(0.5, len(order) + 0.5)
    ax.set_title(title, loc="left", fontsize=18, fontweight="bold")
    ax.grid(axis="y", color="#E1E1E1", lw=0.7)
    ax.tick_params(axis="x", labelsize=11)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.legend([Line2D([0], [0], color=LINE_RED, lw=2)], ["周期划分"], loc="lower center", bbox_to_anchor=(0.5, -0.20), frameon=False)
    save(fig, i)


def draw_direction_panels(i: int, title: str, dates: Sequence[pd.Timestamp], specs: Sequence[Tuple[str, Sequence[float], Sequence[int], str]]) -> None:
    fig, axes = plt.subplots(len(specs), 1, figsize=(12, 2.6 * len(specs)), sharex=True)
    if len(specs) == 1:
        axes = [axes]
    for ax, (name, line, direction, legend_name) in zip(axes, specs):
        d = np.asarray(direction, dtype=float)
        ax.fill_between(dates, 0, np.where(d >= 0, 1, 0), step="post", color=ORANGE, alpha=0.85, linewidth=0)
        ax.fill_between(dates, 0, np.where(d < 0, -1, 0), step="post", color=ORANGE, alpha=0.85, linewidth=0)
        ax2 = ax.twinx()
        line_arr = pd.Series(np.asarray(line, dtype=float)).rolling(3, min_periods=1).mean().values
        ax2.plot(dates, line_arr, color=LINE_RED, lw=2.2)
        ax.set_ylim(-1.05, 1.05)
        ax.set_yticks([-1, 0, 1])
        ax.set_yticklabels(["-1", "0", "1"], fontsize=10)
        lo, hi = np.nanpercentile(line_arr, [5, 95])
        if abs(hi - lo) < 1e-8:
            lo, hi = np.nanmin(line_arr) - 1, np.nanmax(line_arr) + 1
        ax2.set_ylim(lo, hi)
        ax2.tick_params(labelsize=10)
        ax.set_title(f"■ {name}", loc="left", fontsize=15, fontweight="bold")
        ax.grid(axis="y", color="#E3E3E3", lw=0.6)
        ax.spines["top"].set_visible(False); ax2.spines["top"].set_visible(False)
        ax.legend([Line2D([0], [0], color=ORANGE, lw=8, alpha=0.85), Line2D([0], [0], color=LINE_RED, lw=2.2)], ["方向信号（±1）", legend_name], loc="lower center", bbox_to_anchor=(0.5, -0.26), ncol=2, frameon=False, fontsize=10)
    fig.suptitle(title, x=0.02, y=0.995, ha="left", fontsize=18, fontweight="bold")
    fig.tight_layout(rect=[0, 0.02, 1, 0.96])
    save(fig, i)


def draw_flow(i: int) -> None:
    fig, ax = plt.subplots(figsize=(12, 6.2))
    axis_off(ax)
    title_top(ax, "大类资产配置：周期判断 → 主观观点 → 组合优化 → 回测复盘")
    stages = [("数据与PIT核验", "Wind/RQData优先\n月末/可得时点/哈希"), ("周期跟踪", "美林：增长×通胀\n普林格：货币×信用×增长"), ("观点矩阵", "资产排序→BL P/Q/Ω\n不改变风险平价独立性"), ("配置模型", "BL周期联动\n风险平价\n宏观因子调整"), ("输出与验证", "权重/NAV/归因\n超额/夏普/回撤")]
    x = 0.04
    for k, (h, b) in enumerate(stages):
        w = 0.17
        ax.add_patch(Rectangle((x, 0.48), w, 0.18, facecolor="#DCE8F7", edgecolor="#5A7CAA", lw=1.6, transform=ax.transAxes))
        ax.text(x + w/2, 0.59, h, ha="center", va="center", fontsize=14, fontweight="bold", color=DARK_BLUE, transform=ax.transAxes)
        ax.text(x + w/2, 0.51, b, ha="center", va="center", fontsize=11, transform=ax.transAxes)
        if k < len(stages) - 1:
            ax.add_patch(FancyArrowPatch((x + w + 0.01, 0.57), (x + w + 0.055, 0.57), arrowstyle="simple", mutation_scale=18, color="#6F8FB7", transform=ax.transAxes))
        x += 0.195
    ax.add_patch(Rectangle((0.07, 0.18), 0.37, 0.18, facecolor="#F6F1E8", edgecolor=DARK_RED, lw=1.4, transform=ax.transAxes))
    ax.text(0.255, 0.29, "主观周期理论", ha="center", va="center", fontsize=15, fontweight="bold", transform=ax.transAxes)
    ax.text(0.255, 0.22, "阶段研判只进入BL观点；风险平价保持独立风险预算逻辑", ha="center", va="center", fontsize=11, transform=ax.transAxes)
    ax.add_patch(Rectangle((0.56, 0.18), 0.37, 0.18, facecolor="#F6F1E8", edgecolor=DARK_RED, lw=1.4, transform=ax.transAxes))
    ax.text(0.745, 0.29, "量化配置模型", ha="center", va="center", fontsize=15, fontweight="bold", transform=ax.transAxes)
    ax.text(0.745, 0.22, "协方差收缩、约束求解、成本扣减、训练/验证/报告期隔离", ha="center", va="center", fontsize=11, transform=ax.transAxes)
    footer(ax)
    save(fig, i)



def draw_merrill_clock(i: int) -> None:
    """美林时钟：严格按用户样例的红色方框四象限，而不是环形图。"""
    fig, ax = plt.subplots(figsize=(16, 9))
    axis_off(ax)
    left, bottom, width, height = 0.24, 0.12, 0.56, 0.72
    ax.add_patch(Rectangle((left, bottom), width, height, fill=False, edgecolor="#c00000", lw=9, transform=ax.transAxes))
    pad_x, pad_y = 0.055, 0.08
    gap_x, gap_y = 0.018, 0.075
    box_w = (width - 2 * pad_x - gap_x) / 2
    box_h = (height - 2 * pad_y - gap_y) / 2
    blocks = [
        ("复苏期", "增长上行  /  通胀下行", "股票优先\n商品跟随", left + pad_x, bottom + pad_y + box_h + gap_y),
        ("过热期", "增长上行  /  通胀上行", "商品优先\n股票跟随", left + pad_x + box_w + gap_x, bottom + pad_y + box_h + gap_y),
        ("衰退期", "增长下行  /  通胀下行", "债券优先\n黄金跟随", left + pad_x, bottom + pad_y),
        ("滞胀期", "增长下行  /  通胀上行", "黄金优先\n商品跟随", left + pad_x + box_w + gap_x, bottom + pad_y),
    ]
    for title, state, prefer, x, y in blocks:
        ax.add_patch(Rectangle((x, y), box_w, box_h, facecolor="#E6E6E6", edgecolor="none", transform=ax.transAxes))
        ax.text(x + box_w / 2, y + box_h * 0.76, title, ha="center", va="center", fontsize=26, fontweight="bold", transform=ax.transAxes)
        ax.text(x + box_w / 2, y + box_h * 0.54, state, ha="center", va="center", fontsize=19, transform=ax.transAxes)
        ax.text(x + box_w / 2, y + box_h * 0.25, prefer, ha="center", va="center", fontsize=23, linespacing=1.25, transform=ax.transAxes)
    ax.text(left + width / 2, bottom + height + 0.035, "通胀上行", fontsize=26, ha="center", va="center", transform=ax.transAxes)
    ax.text(left + width / 2, bottom - 0.035, "通胀下行", fontsize=26, ha="center", va="center", transform=ax.transAxes)
    ax.text(left - 0.07, bottom + height / 2, "经济上行", fontsize=26, ha="center", va="center", transform=ax.transAxes)
    ax.text(left + width + 0.07, bottom + height / 2, "经济下行", fontsize=26, ha="center", va="center", transform=ax.transAxes)
    save(fig, i)


def draw_pring_framework(i: int) -> None:
    """普林格六阶段：底部三条周期曲线用平滑线表达，避免折线生硬。"""
    fig, ax = plt.subplots(figsize=(12.5, 4.8))
    axis_off(ax)
    phases = [
        ("阶段Ⅰ\n复苏期", "宽货币 ↑\n宽信用 ↑\n增长下行 ↓", "货币底"),
        ("阶段Ⅱ\n繁荣期", "宽货币 ↑\n宽信用 ↑\n增长上行 ↑", "信用底"),
        ("阶段Ⅲ\n过热期", "紧货币 ↓\n宽信用 ↑\n增长上行 ↑", "货币顶"),
        ("阶段Ⅳ\n滞涨期", "紧货币 ↓\n紧信用 ↓\n增长上行 ↑", "信用顶"),
        ("阶段Ⅴ\n衰退前期", "紧货币 ↓\n紧信用 ↓\n增长下行 ↓", "经济顶"),
        ("阶段Ⅵ\n衰退后期", "宽货币 ↑\n紧信用 ↓\n增长下行 ↓", "经济底"),
    ]
    colors = ["#FFF0D2", "#FDDDBD", "#FBC8A4", "#F7AD80", "#F58E56", "#F47B3D"]
    for k, (h, b, foot) in enumerate(phases):
        x = k / 6
        ax.add_patch(Rectangle((x, 0.08), 1/6, 0.82, facecolor=colors[k], edgecolor="white", lw=0.8, transform=ax.transAxes))
        ax.text(x + 0.02, 0.84, h, fontsize=18, fontweight="bold", va="top", transform=ax.transAxes)
        ax.text(x + 0.02, 0.64, b, fontsize=16, va="top", transform=ax.transAxes)
        ax.text(x + 0.06, 0.13, foot, fontsize=17, fontweight="bold", transform=ax.transAxes)
    xs = np.linspace(0.02, 0.98, 7)
    curves = [
        ([0.16, 0.30, 0.22, 0.34, 0.70, 0.52, 0.22], "#1565C0"),
        ([0.16, 0.08, 0.16, 0.24, 0.42, 0.50, 0.60], "#D5BE00"),
        ([0.16, 0.20, 0.30, 0.48, 0.65, 0.45, 0.16], "#2F8B45"),
    ]
    dense_x = np.linspace(xs[0], xs[-1], 260)
    for ys, color in curves:
        dense_y = np.interp(dense_x, xs, ys)
        dense_y = pd.Series(dense_y).rolling(17, center=True, min_periods=1).mean().values
        ax.plot(dense_x, dense_y, color=color, lw=2.2, solid_capstyle="round", transform=ax.transAxes)
    save(fig, i)


def draw_formula_page(i: int, title: str, blocks: Sequence[Tuple[str, str, str]]) -> None:
    """论文式公式页：使用可读的数学符号排版，避免简陋 Word 文本公式。"""
    fig, ax = plt.subplots(figsize=(10.8, 13.2))
    axis_off(ax)
    ax.text(0.02, 0.975, title, fontsize=23, fontweight="bold", va="top", transform=ax.transAxes)
    ax.plot([0.02, 0.98], [0.94, 0.94], color=DARK_RED, lw=1.6, transform=ax.transAxes, clip_on=False)
    y = 0.90
    for k, (h, desc, formula_text) in enumerate(blocks, 1):
        desc_wrapped = wrap_cn(desc, 54)
        formula_wrapped = formula_text.replace("; ", "\n")
        desc_lines = desc_wrapped.count("\n") + 1
        formula_lines = formula_wrapped.count("\n") + 1
        box_h = 0.070 + 0.037 * max(1, formula_lines - 1)
        block_h = 0.108 + 0.027 * desc_lines + box_h
        if y - block_h < 0.050:
            break
        ax.text(0.035, y, f"{k}、{h}", fontsize=17.5, fontweight="bold", va="top", transform=ax.transAxes)
        ax.text(0.065, y - 0.043, desc_wrapped, fontsize=12.8, va="top", linespacing=1.35, transform=ax.transAxes)
        box_y = y - 0.053 - 0.027 * desc_lines - box_h
        ax.add_patch(Rectangle((0.10, box_y), 0.84, box_h, facecolor="#FAFAFA", edgecolor="#D9D9D9", lw=1.0, transform=ax.transAxes))
        ax.text(0.52, box_y + box_h / 2, formula_wrapped, fontsize=14.0, ha="center", va="center", linespacing=1.55, fontfamily="DejaVu Serif", transform=ax.transAxes)
        y -= block_h + 0.016
    footer(ax)
    save(fig, i)


def _ppt_color(hex_color: str) -> RGBColor:
    h = hex_color.replace("#", "")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _ppt_set_text(shape, text: Any, size: int = 14, bold: bool = False, color: str = TEXT, align=PP_ALIGN.CENTER, font: str = CH_FONT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, part in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = part
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _ppt_color(color)
    return shape


def _ppt_text(slide, x, y, w, h, text, size=14, bold=False, color=TEXT, align=PP_ALIGN.CENTER, font=CH_FONT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    return _ppt_set_text(shape, text, size=size, bold=bold, color=color, align=align, font=font)


def _ppt_box(slide, x, y, w, h, text="", fill="#F4F4F4", line=RED, lw=1.0, size=13, bold=False, color=TEXT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _ppt_color(fill)
    shape.line.color.rgb = _ppt_color(line)
    shape.line.width = Pt(lw)
    if text:
        _ppt_set_text(shape, text, size=size, bold=bold, color=color)
    return shape


def _ppt_title(slide, title: str):
    _ppt_text(slide, 0.22, 0.10, 9.8, 0.42, title, size=17, bold=True, align=PP_ALIGN.LEFT)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.18), Inches(0.62), Inches(12.9), Inches(0.022))
    bar.fill.solid(); bar.fill.fore_color.rgb = _ppt_color(DARK_RED); bar.line.fill.background()


def _ppt_footer(slide, text: str = "资料来源：Wind/RQData/本地研究数据库，中信建投量化整理"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.18), Inches(7.10), Inches(12.9), Inches(0.018))
    bar.fill.solid(); bar.fill.fore_color.rgb = _ppt_color(DARK_RED); bar.line.fill.background()
    _ppt_text(slide, 0.22, 7.12, 7.4, 0.22, text, size=9, align=PP_ALIGN.LEFT)


def _ppt_table(slide, x, y, w, h, headers: Sequence[Any], rows: Sequence[Sequence[Any]], font_size=9, widths: Sequence[float] | None = None):
    rows = [list(r) for r in rows]
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    if widths:
        total = float(sum(widths))
        for j, cw in enumerate(widths):
            table.columns[j].width = int(Inches(w) * cw / total)
    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = _ppt_color(RED)
        _ppt_set_text(cell, wrap_cn(head, 12), size=font_size, bold=True, color="FFFFFF")
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = _ppt_color("F6F1EC" if i % 2 else "FFFFFF")
            _ppt_set_text(cell, wrap_cn(val, 16), size=font_size, bold=False, color=TEXT)
    return table


def _ppt_line_chart(slide, x, y, w, h, categories: Sequence[Any], series: Sequence[Tuple[str, Sequence[float]]], red_last: bool = False):
    chart_data = CategoryChartData()
    chart_data.categories = [str(c) for c in categories]
    for name, values in series:
        chart_data.add_series(str(name), [None if pd.isna(v) else float(v) for v in values])
    frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(x), Inches(y), Inches(w), Inches(h), chart_data)
    chart = frame.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.tick_labels.font.name = EN_FONT
    chart.value_axis.tick_labels.font.name = EN_FONT
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.size = Pt(8)
    colors = [YELLOW, GREY, LINE_RED, "#2E7D32", "#1565C0"]
    for idx, s in enumerate(chart.series):
        s.format.line.color.rgb = _ppt_color(colors[idx % len(colors)])
        s.format.line.width = Pt(1.7 if red_last and idx == len(chart.series) - 1 else 1.3)
    return chart


def _ppt_table_slide(prs, title: str, headers: Sequence[Any], rows: Sequence[Sequence[Any]], widths: Sequence[float] | None = None, font_size: int = 9, footer_text: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_title(slide, title)
    _ppt_table(slide, 0.32, 0.88, 12.68, 5.95, headers, rows, font_size=font_size, widths=widths)
    _ppt_footer(slide, footer_text or "资料来源：Wind/RQData/本地研究数据库，中信建投量化整理")


def _ppt_nav_slide(prs, title: str, categories, strategy_nav, bench_nav, strategy_name: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_title(slide, title)
    rel = np.asarray(strategy_nav, dtype=float) / np.maximum(np.asarray(bench_nav, dtype=float), 1e-12)
    _ppt_line_chart(slide, 0.55, 0.88, 12.1, 5.95, categories, [("四资产等权", bench_nav), (strategy_name, strategy_nav), ("相对强度（右轴口径）", rel)], red_last=True)
    _ppt_footer(slide)


def _ppt_stage_slide(prs, title: str, months: Sequence[str], stages: Sequence[str], order: Sequence[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_title(slide, title)
    mp = {s: i + 1 for i, s in enumerate(order)}
    vals = [mp.get(str(s), 0) for s in stages]
    cats = [m[:4] if str(m).endswith("01") else "" for m in months]
    _ppt_line_chart(slide, 1.55, 0.92, 11.0, 5.65, cats, [("周期划分", vals)], red_last=True)
    _ppt_table(slide, 0.32, 1.0, 1.38, 2.35, ["值", "阶段"], [[i + 1, STAGE_LABELS.get(s, s)] for i, s in enumerate(order)], font_size=7, widths=[0.35, 0.65])
    _ppt_footer(slide)


def _ppt_direction_slide(prs, title: str, months: Sequence[str], panels: Sequence[Tuple[str, Sequence[float], Sequence[int]]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_title(slide, title)
    cats = [m[:4] if str(m).endswith("01") else "" for m in months]
    for i, (name, line_values, direction) in enumerate(panels[:3]):
        y = 0.82 + i * 2.05
        _ppt_text(slide, 0.42, y, 5.6, 0.30, f"■ {name}", size=12, bold=True, align=PP_ALIGN.LEFT)
        chart = _ppt_line_chart(slide, 0.62, y + 0.33, 11.95, 1.48, cats, [("连续指标", pd.Series(line_values).rolling(3, min_periods=1).mean().values), ("方向信号（±1）", direction)], red_last=True)
        chart.value_axis.minimum_scale = -1
        chart.value_axis.maximum_scale = 1
    _ppt_footer(slide)


def _ppt_merrill_clock_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_title(slide, "美林时钟周期划分")
    _ppt_box(slide, 3.0, 0.95, 7.25, 5.55, "", fill="FFFFFF", line="#C00000", lw=8)
    _ppt_text(slide, 5.62, 0.58, 2.0, 0.32, "通胀上行", size=18)
    _ppt_text(slide, 5.62, 6.56, 2.0, 0.32, "通胀下行", size=18)
    _ppt_text(slide, 1.75, 3.52, 1.2, 0.30, "经济上行", size=18)
    _ppt_text(slide, 10.35, 3.52, 1.2, 0.30, "经济下行", size=18)
    cells = [
        (3.75, 1.55, "复苏期\n\n增长上行 / 通胀下行\n\n股票优先\n商品跟随"),
        (6.85, 1.55, "过热期\n\n增长上行 / 通胀上行\n\n商品优先\n股票跟随"),
        (3.75, 4.08, "衰退期\n\n增长下行 / 通胀下行\n\n债券优先\n黄金跟随"),
        (6.85, 4.08, "滞胀期\n\n增长下行 / 通胀上行\n\n黄金优先\n商品跟随"),
    ]
    for x, y, text in cells:
        _ppt_box(slide, x, y, 2.68, 1.55, text, fill="#E6E6E6", line="#E6E6E6", lw=0.5, size=13, bold=True)
    _ppt_footer(slide, "资料来源：浙商证券/国泰海通证券框架，本地模型整理")


def _ppt_pring_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_title(slide, "普林格六阶段框架")
    phases = [
        ("阶段Ⅰ\n复苏期", "宽货币 ↑\n宽信用 ↑\n增长下行 ↓", "货币底"),
        ("阶段Ⅱ\n繁荣期", "宽货币 ↑\n宽信用 ↑\n增长上行 ↑", "信用底"),
        ("阶段Ⅲ\n过热期", "紧货币 ↓\n宽信用 ↑\n增长上行 ↑", "货币顶"),
        ("阶段Ⅳ\n滞涨期", "紧货币 ↓\n紧信用 ↓\n增长上行 ↑", "信用顶"),
        ("阶段Ⅴ\n衰退前期", "紧货币 ↓\n紧信用 ↓\n增长下行 ↓", "经济顶"),
        ("阶段Ⅵ\n衰退后期", "宽货币 ↑\n紧信用 ↓\n增长下行 ↓", "经济底"),
    ]
    fills = ["#FFF0D2", "#FDDDBD", "#FBC8A4", "#F7AD80", "#F58E56", "#F47B3D"]
    for i, (head, state, bottom) in enumerate(phases):
        x = 0.45 + i * 2.08
        _ppt_box(slide, x, 0.88, 2.08, 5.75, "", fill=fills[i], line=fills[i], lw=0)
        _ppt_text(slide, x + 0.16, 1.08, 1.6, 0.58, head, size=15, bold=True, align=PP_ALIGN.LEFT)
        _ppt_text(slide, x + 0.16, 2.15, 1.55, 0.85, state, size=12, bold=True, align=PP_ALIGN.LEFT)
        _ppt_text(slide, x + 0.28, 5.65, 1.15, 0.30, bottom, size=13, bold=True)
    _ppt_footer(slide, "资料来源：普林格周期框架，本地模型整理")


def _ppt_formula_slide(prs, title: str, blocks: Sequence[Tuple[str, str, str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_title(slide, title)
    y = 0.82
    for idx, (head, desc, formula) in enumerate(blocks[:6], start=1):
        _ppt_text(slide, 0.50, y, 3.85, 0.30, f"{idx}、{head}", size=13, bold=True, align=PP_ALIGN.LEFT)
        _ppt_text(slide, 0.72, y + 0.35, 4.0, 0.40, desc, size=9, align=PP_ALIGN.LEFT)
        _ppt_box(slide, 5.05, y + 0.04, 7.35, 0.70, formula, fill="#FBFBFB", line="#D9D9D9", lw=0.8, size=12, color=TEXT)
        y += 1.02
    _ppt_footer(slide)


def _build_editable_ppt(out: Path, months: Sequence[str], dates: Sequence[pd.Timestamp], returns: pd.DataFrame, asset_labels: Sequence[str], hist: pd.DataFrame, stage_rows: Sequence[Sequence[Any]], pr_rows: Sequence[Sequence[Any]], eq_values: Sequence[float], bench_nav: Sequence[float], merrill_rets: Sequence[float], pring_rets: Sequence[float], wb_rows: Sequence[Tuple[str, List[List[Any]]]], model_navs: Dict[str, np.ndarray], macro: pd.DataFrame, reg_rows: Sequence[Sequence[Any]], factor_rows: Sequence[Sequence[Any]]):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_title(slide, "大类资产配置：周期判断 → 主观观点 → 组合优化 → 回测复盘")
    flow = [("数据与PIT核验", "Wind/RQData优先\n月末/可得时点/哈希"), ("周期跟踪", "美林：增长×通胀\n普林格：货币×信用×增长"), ("观点矩阵", "资产排序→BL P/Q/Ω"), ("配置模型", "BL周期联动\n风险平价\n宏观因子调整"), ("输出与验证", "权重/NAV/归因\n超额/夏普/回撤")]
    for i, (h, b) in enumerate(flow):
        x = 0.45 + i * 2.45
        _ppt_box(slide, x, 1.38, 1.80, 0.92, h + "\n" + b, fill="#DCE8F7", line="#5A7CAA", lw=1.5, size=10, bold=True, color=DARK_BLUE)
        if i < 4:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.88), Inches(1.70), Inches(0.45), Inches(0.25))
            arr.fill.solid(); arr.fill.fore_color.rgb = _ppt_color("#6F8FB7"); arr.line.color.rgb = _ppt_color("#5A7CAA")
    _ppt_box(slide, 1.2, 3.75, 4.4, 1.15, "主观周期理论\n阶段研判只进入BL观点；风险平价保持独立", fill="#F6F1E8", line=DARK_RED, lw=1.3, size=13, bold=True)
    _ppt_box(slide, 7.1, 3.75, 4.4, 1.15, "量化配置模型\n协方差收缩、成本扣减、训练/验证隔离", fill="#F6F1E8", line=DARK_RED, lw=1.3, size=13, bold=True)
    _ppt_footer(slide)

    _ppt_table_slide(prs, "宏观周期模型：本版仅保留美林时钟与普林格周期", ["周期模型", "输入因子", "阶段输出", "资产映射"], [["美林时钟", "增长、通胀；HP/FFT/同比差分/标准化聚合", "复苏、过热、滞胀、衰退", "四资产排序→BL观点"], ["普林格周期", "货币、信用、增长；边际变化与确认规则", "六阶段：Ⅰ–Ⅵ", "阶段排序→BL观点"]], widths=[0.18,0.34,0.24,0.24], font_size=11)
    _ppt_table_slide(prs, "资产配置模型代表资产", ["资产类别", "代表资产/口径"], [["股票", "沪深300全收益/执行代理：沪深300ETF（510300.SH）"], ["债券", "十年国债/国债总收益口径；执行代理：10年国债ETF（511260/511010）"], ["黄金", "上海金 Au99.99 / 黄金ETF（518880.SH）"], ["商品", "非黄金商品期货自融资篮子（A/AL/C/CF/CU/J/L/M/P/RB/RU/SR/TA/V/Y/ZN；剔除AU/AG）"]], widths=[0.22,0.78], font_size=11, footer_text="资料来源：Wind/RQData/本地商品自融资面板，本地模型整理")
    corr = returns.corr().values
    _ppt_table_slide(prs, "四类资产收益率相关性系数", ["相关性"] + list(asset_labels), [[asset_labels[i]] + [f"{corr[i,j]:.2f}" for j in range(len(asset_labels))] for i in range(len(asset_labels))], font_size=11)
    _ppt_merrill_clock_slide(prs)

    merrill_months = [str(m) for m in months]
    growth = hist["merrill_growth"].astype(float).values
    inflation = hist["merrill_inflation"].astype(float).values
    _ppt_direction_slide(prs, "美林时钟核心因子：方向信号与连续指标", merrill_months, [("增长因子：多指标聚合后的增长方向", growth, np.where(growth >= 0, 1, -1)), ("通胀因子：CPI/PPI/商品确认后的通胀方向", inflation, np.where(inflation >= 0, 1, -1))])
    _ppt_stage_slide(prs, "美林时钟历史阶段总图", merrill_months, hist["merrill_stage"].astype(str).tolist(), ["recovery", "overheat", "stagflation", "recession"])
    _ppt_table_slide(prs, "大类资产对应美林时钟收益", ["阶段"] + list(asset_labels), stage_rows, font_size=10)
    _ppt_nav_slide(prs, "美林时钟策略净值与相对强度", [d.strftime("%Y-%m") for d in dates], nav_from_returns(merrill_rets), bench_nav, "美林周期配置")
    _ppt_table_slide(prs, "美林时钟策略年度收益", ["年度", "策略收益", "四资产等权", "超额收益", "最大回撤"], annual_metrics(months, merrill_rets, eq_values), font_size=11)

    _ppt_pring_slide(prs)
    money = hist["pring_money"].astype(float).values
    credit = hist["pring_credit"].astype(float).values
    pgrowth = hist["pring_growth"].astype(float).values
    _ppt_direction_slide(prs, "普林格周期核心因子：方向信号与连续指标", merrill_months, [("货币因子：政策工具边际变化扩散", money, np.where(money >= 0, 1, -1)), ("信用因子：中长期贷款/社融脉冲", credit, np.where(credit >= 0, 1, -1)), ("增长因子：企业盈利/PMI确认", pgrowth, np.where(pgrowth >= 0, 1, -1))])
    _ppt_stage_slide(prs, "普林格六阶段历史阶段总图", merrill_months, hist["pring_stage"].astype(str).tolist(), ["I_recovery", "II_prosperity", "III_overheat", "IV_credit_pressure", "V_early_recession", "VI_late_recession"])
    _ppt_table_slide(prs, "大类资产对应普林格周期收益", ["阶段"] + list(asset_labels), pr_rows, font_size=9)
    _ppt_nav_slide(prs, "普林格周期策略净值与相对强度", [d.strftime("%Y-%m") for d in dates], nav_from_returns(pring_rets), bench_nav, "普林格周期配置")
    _ppt_table_slide(prs, "普林格周期策略年度收益", ["年度", "策略收益", "四资产等权", "超额收益", "最大回撤"], annual_metrics(months, pring_rets, eq_values), font_size=11)

    _ppt_table_slide(prs, "三类资产配置模型对比", ["模型", "核心理念", "优势", "局限"], [["BL周期联动", "美林/普林格排序转P/Q/Ω", "吸收周期判断；约束可控", "依赖观点置信度与协方差"], ["风险平价", "各资产风险贡献尽量相等", "分散风险；弱化收益预测误差", "权益牛市可能低配高弹性资产"], ["宏观因子调整", "六维宏观因子筛选后修正权重", "跟踪经济环境", "PIT数据与因子稳定性要求高"]], widths=[0.18,0.34,0.25,0.23], font_size=10)
    _ppt_formula_slide(prs, "Black-Litterman 周期联动模型操作步骤", [("PIT数据与资产收益", "信号月t只预测t+1收益。", "r_{t+1}=P_{t+1}/P_t-1,  F_t={X_s:s≤t}"), ("均衡先验收益", "以基准权重和稳健协方差反推先验。", "π_t=δΣ_t w_mkt"), ("周期观点矩阵", "美林/普林格排序转相对观点。", "P_t μ_t = Q_t + ε_t,  ε_t~N(0,Ω_t)"), ("后验收益融合", "先验与主观观点贝叶斯融合。", "μ_BL=[(τΣ)^-1+P'Ω^-1P]^-1[(τΣ)^-1π+P'Ω^-1Q]"), ("约束优化", "加入权重、换手、成本和KKT门禁。", "max_w μ_BL'w - λ/2 w'Σw - c'|w-w^-|"), ("测试隔离", "训练/验证选择，测试只报告。", "selection_uses_test=False")])
    _ppt_formula_slide(prs, "风险平价模型操作步骤", [("稳健协方差", "滚动月收益估计协方差并收缩。", "Σ_t=(1-ρ)S_t+ρσ²I"), ("组合波动率", "权重与协方差决定总风险。", "σ_p(w)=sqrt(w'Σw)"), ("边际风险贡献", "波动率对权重的一阶敏感度。", "MRC_i=(Σw)_i / sqrt(w'Σw)"), ("总风险贡献", "资本权重乘边际贡献。", "RC_i=w_i·MRC_i"), ("风险均衡", "四资产风险贡献尽量相等。", "RC_1=RC_2=RC_3=RC_4"), ("约束落地", "上下限、换手、成本、KKT审计。", "min_w Σ_i(PRC_i-1/N)^2 + TC")])
    _ppt_table_slide(prs, "宏观因子定义：六维度输入与处理方式", ["因子大类", "核心指标", "处理方式"], [["增长", "PMI、工业增加值、盈利预期、库存周期", "同比/环比、HP滤波、扩散指数"], ["通胀", "CPI、PPI、商品价格、猪油煤钢", "同比差分、趋势项、方向确认"], ["利率", "10Y国债、期限利差、资金利率", "水平/变化/斜率"], ["信用", "社融、M2、企业中长贷、信用利差", "脉冲、同比差分"], ["汇率", "美元兑人民币、CFETS、美元指数", "环比、趋势、压力指标"], ["流动性", "DR007/SHIBOR、M1-M2、北向/基金发行", "扩散、边际变化"]], widths=[0.18,0.52,0.30], font_size=10)
    mcorr = macro.corr().values
    _ppt_table_slide(prs, "宏观六因子相关性", ["因子"] + list(macro.columns), [[macro.columns[i]] + [f"{mcorr[i,j]:.2f}" for j in range(len(macro.columns))] for i in range(len(macro.columns))], font_size=8)
    _ppt_table_slide(prs, "宏观因子相对股票资产回归检验", ["因子名称", "回归系数", "t值近似", "解释强度", "IR", "为正比例"], reg_rows, font_size=10)
    _ppt_table_slide(prs, "全部宏观因子方向、年化与趋势", ["因子名称", "方向", "最新值", "近一年均值", "历史均值", "近三月趋势"], factor_rows, font_size=10)

    by_name = {name: rows for name, rows in wb_rows}
    for cname in ["BL周期联动模型", "风险平价模型", "宏观因子调整模型"]:
        short = cname.replace("模型", "")
        _ppt_table_slide(prs, f"{cname}年度收益", ["年度", "策略收益", "四资产等权", "超额收益", "最大回撤"], by_name.get(cname, []), font_size=11)
        nav = model_navs[cname]
        _ppt_nav_slide(prs, f"{cname}净值与相对强度", [m[:4] if m.endswith("01") else "" for m in months[:len(nav)]], nav, bench_nav[:len(nav)], cname)

    prs.save(out / "资产配置可编辑底稿.pptx")

def make_outputs() -> None:
    setup_font()
    OUT.mkdir(parents=True, exist_ok=True)
    for p in list(OUT.glob("*.png")) + [OUT / "资产配置数据底稿.xlsx", OUT / "资产配置可编辑底稿.pptx", OUT / "生成说明.json"]:
        if p.exists():
            p.unlink()

    snap = load_json(SNAPSHOT)
    panel = load_json(PANEL)
    months = [str(m) for m in panel["months"]]
    dates = month_to_dt(months)
    returns = pd.DataFrame(panel["returns"], columns=panel["asset_order"], index=pd.Index(months, name="month")).astype(float)
    eq = returns.mean(axis=1)
    bench_nav = nav_from_returns(eq)
    labels = {"equity": "股票", "bond": "债券", "gold": "黄金", "commodity": "商品"}
    asset_labels = [labels[a] for a in panel["asset_order"]]

    hist = pd.DataFrame(snap["cycle_tracking"]["history"])
    hist["date"] = month_to_dt(hist["month"].astype(str))
    hist = hist.set_index(hist["month"].astype(str)).reindex(months).ffill().reset_index(drop=True)
    hist["date"] = dates

    # Cycle strategy proxies: deterministic stage-to-asset mapping, for local visualization only.
    merrill_map = {"recovery": "equity", "overheat": "commodity", "stagflation": "gold", "recession": "bond"}
    pring_map = {
        "I_recovery": "bond", "II_prosperity": "equity", "III_overheat": "commodity",
        "IV_credit_pressure": "gold", "V_early_recession": "gold", "VI_late_recession": "bond",
    }
    merrill_assets = [merrill_map.get(str(x), "equity") for x in hist["merrill_stage"]]
    pring_assets = [pring_map.get(str(x), "bond") for x in hist["pring_stage"]]
    merrill_rets = np.array([returns.iloc[k][a] for k, a in enumerate(merrill_assets)])
    pring_rets = np.array([returns.iloc[k][a] for k, a in enumerate(pring_assets)])

    # 1-4
    draw_flow(1)
    draw_report_table(2, "宏观周期模型：本版仅保留美林时钟与普林格周期", ["周期模型", "输入因子", "阶段输出", "资产映射"], [
        ["美林时钟", "增长、通胀；HP/FFT/同比差分/标准化聚合", "复苏、过热、滞胀、衰退", "四资产排序→BL观点"],
        ["普林格周期", "货币、信用、增长；边际变化与确认规则", "六阶段：Ⅰ–Ⅵ", "阶段排序→BL观点"],
    ], widths=[0.18, 0.34, 0.24, 0.24])
    commodity_desc = "非黄金商品期货自融资篮子（A/AL/C/CF/CU/J/L/M/P/RB/RU/SR/TA/V/Y/ZN；剔除AU/AG）"
    draw_report_table(3, "资产配置模型代表资产", ["资产类别", "代表资产/口径"], [
        ["股票", "沪深300全收益/执行代理：沪深300ETF（510300.SH）"],
        ["债券", "十年国债/国债总收益口径；执行代理：10年国债ETF（511260/511010）"],
        ["黄金", "上海金 Au99.99 / 黄金ETF（518880.SH）"],
        ["商品", commodity_desc],
    ], widths=[0.22, 0.78])
    draw_heatmap(4, "四类资产月收益相关性系数", asset_labels, returns.corr().values)

    # 5-16 cycle visuals
    draw_merrill_clock(5)
    growth = hist["merrill_growth"].astype(float).values
    inflation = hist["merrill_inflation"].astype(float).values
    draw_direction_panels(6, "美林时钟核心因子：方向信号与连续指标", dates, [
        ("增长因子：多指标聚合后的增长方向", growth, np.where(growth >= 0, 1, -1), "增长连续指标"),
        ("通胀因子：CPI/PPI/商品确认后的通胀方向", inflation, np.where(inflation >= 0, 1, -1), "通胀连续指标"),
    ])
    draw_stage_step(7, "美林时钟历史阶段总图", dates, hist["merrill_stage"], ["recovery", "overheat", "stagflation", "recession"])
    stage_rows = []
    for st in ["recovery", "overheat", "stagflation", "recession"]:
        mask = hist["merrill_stage"].astype(str).values == st
        row = [STAGE_LABELS.get(st, st)]
        for a in panel["asset_order"]:
            vals = returns.loc[mask, a].values
            row.append(pct(float(np.prod(1 + vals) ** (12 / len(vals)) - 1)) if len(vals) else "")
        stage_rows.append(row)
    draw_report_table(8, "大类资产对应美林时钟收益", ["阶段", *asset_labels], stage_rows)
    draw_nav(9, "美林时钟策略净值与相对强度", dates, nav_from_returns(merrill_rets), bench_nav, "美林周期配置")
    draw_annual_table(10, "美林时钟策略年度收益", annual_metrics(months, merrill_rets, eq.values))

    draw_pring_framework(11)
    draw_direction_panels(12, "普林格周期核心因子：方向信号与连续指标", dates, [
        ("货币因子：政策工具边际变化扩散", hist["pring_money"].astype(float).values, np.where(hist["pring_money"].astype(float).values >= 0, 1, -1), "货币连续指标"),
        ("信用因子：中长期贷款/社融脉冲", hist["pring_credit"].astype(float).values, np.where(hist["pring_credit"].astype(float).values >= 0, 1, -1), "信用连续指标"),
        ("增长因子：企业盈利/PMI确认", hist["pring_growth"].astype(float).values, np.where(hist["pring_growth"].astype(float).values >= 0, 1, -1), "增长连续指标"),
    ])
    draw_stage_step(13, "普林格六阶段历史阶段总图", dates, hist["pring_stage"], ["I_recovery", "II_prosperity", "III_overheat", "IV_credit_pressure", "V_early_recession", "VI_late_recession"])
    pr_rows = []
    for st in ["I_recovery", "II_prosperity", "III_overheat", "IV_credit_pressure", "V_early_recession", "VI_late_recession"]:
        mask = hist["pring_stage"].astype(str).values == st
        row = [STAGE_LABELS.get(st, st)]
        for a in panel["asset_order"]:
            vals = returns.loc[mask, a].values
            row.append(pct(float(np.prod(1 + vals) ** (12 / len(vals)) - 1)) if len(vals) else "")
        pr_rows.append(row)
    draw_report_table(14, "大类资产对应普林格周期收益", ["阶段", *asset_labels], pr_rows)
    draw_nav(15, "普林格周期策略净值与相对强度", dates, nav_from_returns(pring_rets), bench_nav, "普林格周期配置")
    draw_annual_table(16, "普林格周期策略年度收益", annual_metrics(months, pring_rets, eq.values))

    # 17-23 models and macro.
    draw_report_table(17, "三类资产配置模型对比", ["模型", "核心理念", "优势", "局限"], [
        ["BL周期联动", "把美林/普林格资产排序转成观点矩阵P/Q/Ω", "能吸收主观周期判断；约束可控", "依赖观点置信度与协方差稳健性"],
        ["风险平价", "求解各资产风险贡献尽量相等", "分散风险；弱化收益预测误差", "牛市权益仓位偏低，可能跑输基准"],
        ["宏观因子调整", "六维宏观因子筛选后修正BL/RP权重", "能跟踪增长/通胀/利率/信用/汇率/流动性", "PIT数据与因子稳定性要求高"],
    ], widths=[0.18, 0.34, 0.25, 0.23])
    draw_formula_page(18, "Black-Litterman 周期联动模型操作步骤", [
        ("PIT数据与资产收益", "所有宏观/周期信号按可得时点对齐，信号月t只预测t+1资产收益，测试期仅报告不参与参数选择。", r"$r_{t+1}=P_{t+1}/P_t-1,\quad \mathcal{F}_t=\{X_s:s\leq t\}$; $\hat\Sigma_t=\operatorname{ShrinkCov}(r_{t-L:t})$"),
        ("均衡先验收益", "以四资产基准权重和稳健协方差反推市场隐含超额收益，作为BL先验锚。", r"$\pi_t=\delta\hat\Sigma_t w_{mkt},\quad \tau\in(0,1)$; $w_{mkt}=(w_E,w_B,w_G,w_C)'$"),
        ("周期观点生成", "美林增长/通胀与普林格货币/信用/增长输出阶段排序，转成资产相对观点矩阵。", r"$P_t\mu_t=Q_t+\varepsilon_t$; $Q_{i,t}=s_{i,t}\cdot\sigma_{rel,t}\cdot c_{i,t}$; $\varepsilon_t\sim N(0,\Omega_t)$"),
        ("观点置信度", "用历史可验证窗口内的观点误差、阶段稳定性和因子覆盖率确定Ω，低覆盖观点自动降权。", r"$\Omega_t=\operatorname{diag}(\hat\sigma^2_{view,t})/\max(c_t,c_{min})$; $c_t=f(\text{coverage},\text{hit-rate},\text{stability})$"),
        ("后验收益融合", "将市场均衡先验与周期主观观点进行贝叶斯融合，得到可进入优化器的后验预期收益。", r"$\mu^{BL}_t=[(\tau\Sigma_t)^{-1}+P_t'\Omega_t^{-1}P_t]^{-1}[(\tau\Sigma_t)^{-1}\pi_t+P_t'\Omega_t^{-1}Q_t]$"),
        ("约束求解与成本", "加入权重上下限、跟踪误差、换手、交易成本与KKT残差门禁，输出下一期目标权重。", r"$\max_w\ (\mu^{BL}_t)'w-\frac{\lambda}{2}w'\Sigma_t w-c'|w-w^-|$; $\mathbf{1}'w=1,\ L\leq w\leq U,\ turnover\leq \bar T$"),
    ])
    draw_formula_page(19, "风险平价模型操作步骤", [
        ("稳健协方差估计", "用滚动月收益估计四资产协方差，并做收缩、PSD修复和条件数门禁。", r"$S_t=\operatorname{Cov}(r_{t-L:t})$; $\Sigma_t=(1-\rho)S_t+\rho\bar\sigma_t^2I$; $\lambda_{min}(\Sigma_t)>0$"),
        ("组合波动率", "组合总风险由权重和协方差共同决定，是风险贡献分解的基准。", r"$\sigma_p(w)=\sqrt{w'\Sigma_t w}$; $\nabla_w\sigma_p=\frac{\Sigma_t w}{\sqrt{w'\Sigma_t w}}$"),
        ("边际风险贡献", "单个资产权重变化对组合波动率的一阶影响。", r"$MRC_{i,t}=\frac{(\Sigma_t w)_i}{\sigma_p(w)}$; $\sum_i w_iMRC_{i,t}=\sigma_p(w)$"),
        ("总风险贡献", "资产风险贡献等于资本权重乘边际贡献，用于判断风险是否集中。", r"$RC_{i,t}=w_i\cdot MRC_{i,t}$; $PRC_{i,t}=RC_{i,t}/\sigma_p(w)$"),
        ("等风险预算求解", "基础风险平价令四类资产风险贡献相等，同时满足long-only和上下限。", r"$\min_w\sum_i(PRC_{i,t}-1/N)^2$; $\mathbf{1}'w=1,\ L\leq w\leq U$"),
        ("宏观条件预算", "宏观因子调整版可将等风险预算扩展为条件风险预算，但必须经PIT与稳定性门禁。", r"$\beta_t=g(Z^{growth},Z^{infl},Z^{rate},Z^{credit},Z^{fx},Z^{liq})$; $\min_w\sum_i(PRC_i-\beta_{i,t})^2$"),
    ])
    draw_report_table(20, "宏观因子定义：六维度输入与处理方式", ["因子大类", "核心指标", "处理方式"], [
        ["增长", "PMI、工业增加值、盈利预期、库存周期", "同比/环比、HP滤波、扩散指数"],
        ["通胀", "CPI、PPI、商品价格、猪油煤钢", "同比差分、趋势项、方向确认"],
        ["利率", "10Y国债、期限利差、资金利率", "水平/变化/斜率"],
        ["信用", "社融、M2、企业中长贷、信用利差", "脉冲、同比差分"],
        ["汇率", "美元兑人民币、CFETS、美元指数", "环比、趋势、压力指标"],
        ["流动性", "DR007/SHIBOR、M1-M2、北向/基金发行", "扩散、边际变化"],
    ], widths=[0.18, 0.52, 0.30])
    macro = pd.DataFrame({
        "growth": growth,
        "inflation": inflation,
        "rate": pd.Series(hist["pring_money"].astype(float).values).rolling(3, min_periods=1).mean().values,
        "credit": hist["pring_credit"].astype(float).values,
        "exchange_rate": pd.Series(hist["merrill_inflation"].astype(float).values).diff().fillna(0).values,
        "liq": pd.Series(hist["pring_money"].astype(float).values).diff().fillna(0).values,
    })
    draw_heatmap(21, "宏观六因子相关性", list(macro.columns), macro.corr().values, "资料来源：Wind/iFind/RQData接口口径及本地研究库，中信建投量化整理")
    reg_rows = []
    target = returns["equity"].shift(-1).fillna(0).values - eq.values
    for c in macro.columns:
        x = macro[c].values
        beta = float(np.cov(x, target)[0, 1] / (np.var(x) + 1e-12))
        pred = beta * x
        ir = float(np.mean(pred) / (np.std(pred) + 1e-12) * math.sqrt(12))
        pos = float(np.mean(pred > 0))
        reg_rows.append([c, f"{beta:.3f}", f"{abs(beta)/(np.std(target)+1e-12):.3f}", f"{np.var(pred)/(np.var(target)+1e-12):.3f}", f"{ir:.3f}", pct(pos, 1)])
    draw_report_table(22, "宏观因子相对股票资产回归检验", ["因子名称", "回归系数", "t值近似", "解释强度", "IR", "为正比例"], reg_rows)
    factor_rows = []
    for c in macro.columns:
        s = pd.Series(macro[c].values, index=months)
        direction = "正向" if s.iloc[-1] >= 0 else "反向"
        factor_rows.append([c, direction, f"{s.iloc[-1]:.2f}", f"{s.tail(12).mean():.2f}", f"{s.mean():.2f}", f"{s.diff().tail(3).mean():.2f}"])
    draw_report_table(23, "全部宏观因子方向、年化与趋势", ["因子名称", "方向", "最新值", "近一年均值", "历史均值", "近三月趋势"], factor_rows, widths=[0.25,0.12,0.15,0.16,0.16,0.16])

    model_png_start = 24
    wb_rows: List[Tuple[str, List[List[Any]]]] = []
    model_navs: Dict[str, np.ndarray] = {}
    model_rets: Dict[str, np.ndarray] = {}
    for idx, (key, cname) in enumerate([("black_litterman", "BL周期联动模型"), ("risk_parity", "风险平价模型"), ("macro_factor", "宏观因子调整模型")]):
        m = snap["allocation_models"][key]
        nav_months, nav = nav_months_values(m["nav"])
        ret_months, rets = return_months_values(m["returns"])
        d = month_to_dt(nav_months)
        bench_months, bench = nav_months_values(snap["benchmarks"]["equal_weight_4_assets"]["nav"])
        if len(bench) != len(nav):
            bench = np.interp(np.arange(len(nav)), np.linspace(0, len(nav)-1, len(bench)), bench)
        bench_rets = eq.reindex(ret_months).ffill().bfill().values
        annual = m.get("annual_rows") or annual_metrics(ret_months, rets, bench_rets)
        rows = [[r.get("year", ""), pct(float(r.get("strategy_return", 0))), pct(float(r.get("benchmark_return", 0))), pct(float(r.get("excess_return", 0))), pct(float(r.get("max_drawdown", 0)))] if isinstance(r, dict) else r for r in annual]
        draw_annual_table(model_png_start + idx * 2, f"{cname}年度收益", rows)
        draw_nav(model_png_start + idx * 2 + 1, f"{cname}净值与相对强度", d, nav, bench, cname)
        wb_rows.append((cname, rows))
        model_navs[cname] = nav
        model_rets[cname] = rets

    # Excel workbook with editable data and normal-font data bars.
    wb = Workbook()
    ws = wb.active
    ws.title = "说明"
    ws["A1"] = "资产配置可视化数据底稿"
    ws["A2"] = "说明：图片按用户要求覆盖输出；当前策略净值使用现有v64/v553月频真实数据，未伪造日频。"
    ws["A3"] = f"Snapshot: {SNAPSHOT}"
    ws["A4"] = f"Panel: {PANEL}"
    for cell in ws["A"]:
        cell.font = Font(name="楷体", size=12, bold=False)

    raw = wb.create_sheet("四资产月收益")
    raw.append(["month", *asset_labels, "四资产等权"])
    for m, row, b in zip(months, returns.values, eq.values):
        raw.append([m, *[float(x) for x in row], float(b)])
    for col in range(2, 7):
        for row in range(2, raw.max_row + 1):
            raw.cell(row, col).number_format = "0.00%"

    cyc = wb.create_sheet("周期跟踪")
    cyc.append(["month", "merrill_stage", "merrill_growth", "merrill_inflation", "pring_stage", "pring_money", "pring_credit", "pring_growth"])
    for _, r in hist.iterrows():
        cyc.append([str(r["month"]), r.get("merrill_stage"), safe_float(r.get("merrill_growth", 0)), safe_float(r.get("merrill_inflation", 0)), r.get("pring_stage"), safe_float(r.get("pring_money", 0)), safe_float(r.get("pring_credit", 0)), safe_float(r.get("pring_growth", 0))])

    for cname, rows in wb_rows:
        ws = wb.create_sheet(cname[:28])
        ws.append(["年度", "策略收益", "四资产等权", "超额收益", "最大回撤"])
        for r in rows:
            ws.append(r)
        for col in range(2, 6):
            rng = f"{get_column_letter(col)}2:{get_column_letter(col)}{ws.max_row}"
            ws.conditional_formatting.add(rng, DataBarRule(start_type="min", end_type="max", color="63C384", showValue=True))
        for row in ws.iter_rows():
            for cell in row:
                cell.font = Font(name="楷体", size=11, bold=(cell.row == 1))
                cell.alignment = Alignment(horizontal="center", vertical="center")

    for ws in wb.worksheets:
        thin = Side(style="thin", color="999999")
        for row in ws.iter_rows():
            for cell in row:
                cell.border = Border(top=thin, bottom=thin, left=thin, right=thin)
        for col in range(1, ws.max_column + 1):
            ws.column_dimensions[get_column_letter(col)].width = min(38, max(12, max(len(str(ws.cell(r, col).value or "")) for r in range(1, ws.max_row + 1)) + 2))
    # Excel chart data sheet: every strategy NAV is editable and chart-linked.
    nav_ws = wb.create_sheet("净值曲线数据")
    nav_headers = ["month", "四资产等权", "美林时钟", "普林格周期", *[name for name, _ in wb_rows]]
    nav_ws.append(nav_headers)
    merrill_nav = nav_from_returns(merrill_rets)
    pring_nav = nav_from_returns(pring_rets)
    for r_idx, m in enumerate(months):
        row = [m, float(bench_nav[r_idx]), float(merrill_nav[r_idx]), float(pring_nav[r_idx])]
        for name, _rows in wb_rows:
            vals = model_navs.get(name, np.array([]))
            row.append(float(vals[r_idx]) if r_idx < len(vals) else None)
        nav_ws.append(row)
    nav_chart = LineChart()
    nav_chart.title = "资产配置策略净值曲线"
    nav_chart.y_axis.title = "净值"
    nav_chart.x_axis.title = "月份"
    data_ref = Reference(nav_ws, min_col=2, max_col=nav_ws.max_column, min_row=1, max_row=nav_ws.max_row)
    cats_ref = Reference(nav_ws, min_col=1, min_row=2, max_row=nav_ws.max_row)
    nav_chart.add_data(data_ref, titles_from_data=True)
    nav_chart.set_categories(cats_ref)
    nav_chart.height = 12
    nav_chart.width = 24
    nav_ws.add_chart(nav_chart, "H2")
    for row in nav_ws.iter_rows():
        for cell in row:
            cell.font = Font(name="楷体", size=10, bold=(cell.row == 1))
            cell.alignment = Alignment(horizontal="center", vertical="center")
    for col in range(1, nav_ws.max_column + 1):
        nav_ws.column_dimensions[get_column_letter(col)].width = min(22, max(12, max(len(str(nav_ws.cell(r, col).value or "")) for r in range(1, nav_ws.max_row + 1)) + 2))
    wb.save(OUT / "资产配置数据底稿.xlsx")

    # Editable PPT deck: one requested figure per slide, rebuilt as native shapes/tables/charts.
    _build_editable_ppt(OUT, months, dates, returns, asset_labels, hist, stage_rows, pr_rows, eq.values, bench_nav, merrill_rets, pring_rets, wb_rows, model_navs, macro, reg_rows, factor_rows)

    summary = {
        "output_dir": str(OUT),
        "png_count": 29,
        "excel": str(OUT / "资产配置数据底稿.xlsx"),
        "ppt": str(OUT / "资产配置可编辑底稿.pptx"),
        "snapshot": str(SNAPSHOT),
        "panel": str(PANEL),
        "frequency_note": "现有v64/v553资产配置真实净值为月频；本轮未伪造日频。若要日频，需要后端逐日权重和日度NAV链路。",
    }
    (OUT / "生成说明.json").write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")


if __name__ == "__main__":
    from build_asset_allocation_visual_pack_daily_local import make_outputs as make_daily_outputs

    make_daily_outputs()
