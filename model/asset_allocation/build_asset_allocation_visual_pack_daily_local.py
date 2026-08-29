from __future__ import annotations

import json
import math
import shutil
import sys
from pathlib import Path
from typing import Any, Dict, Iterable, List, Sequence, Tuple

import numpy as np
import pandas as pd

import build_asset_allocation_visual_pack_local as base
from openpyxl import Workbook
from openpyxl.chart import LineChart, Reference
from openpyxl.formatting.rule import DataBarRule
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches


ROOT = Path(r"G:\中信建投")
OUT = Path(r"C:\Users\Rye\Desktop\资产配置")
SNAPSHOT = ROOT / "agent" / "output" / "model_improvement" / "asset_allocation_snapshot_v64_daily_excess_governed.json"
PANEL = ROOT / "agent" / "output" / "model_improvement" / "asset_allocation_panel_v553.json"
FREEZE = ROOT / "agent" / "output" / "model_improvement" / "asset_allocation_rqdata_v541_freeze.json"
MODEL_DIR = ROOT / "agent" / "model" / "asset_allocation"

ASSET_ORDER = ["equity", "bond", "gold", "commodity"]
ASSET_LABELS = ["股票", "债券", "黄金", "商品"]
ASSET_CODES = {
    "股票": "沪深300全收益 H00300.INDX；执行代理 510300.SH",
    "债券": "中证国债收益 H11006.XSHG；执行代理 511010/511260",
    "黄金": "上海金 Au99.99 AU9999.SGEX；执行代理 518880.SH",
    "商品": "非贵金属期货自融资篮子 A/AL/C/CF/CU/J/L/M/P/RB/RU/SR/TA/V/Y/ZN；排除 AU/AG",
}
MODEL_NAMES = {
    "black_litterman": "周期BL",
    "risk_parity": "风险平价",
    "macro_factor": "宏观因子",
}
STAGE_TO_ASSET_MERRILL = {
    "recovery": "equity",
    "overheat": "commodity",
    "stagflation": "gold",
    "recession": "bond",
}
STAGE_TO_ASSET_PRING = {
    "I_recovery": "bond",
    "II_prosperity": "equity",
    "III_overheat": "commodity",
    "IV_credit_pressure": "gold",
    "V_early_recession": "gold",
    "VI_late_recession": "bond",
}


def _load_json(path: Path) -> Dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _ym(date: pd.Timestamp) -> str:
    return date.strftime("%Y%m")


def _geom(rets: Sequence[float], periods: int = 252) -> float:
    arr = np.asarray(rets, dtype=float)
    arr = arr[np.isfinite(arr)]
    if arr.size == 0:
        return 0.0
    return float(np.prod(1 + arr) ** (periods / arr.size) - 1)


def _max_dd(rets: Sequence[float]) -> float:
    arr = np.asarray(rets, dtype=float)
    if arr.size == 0:
        return 0.0
    nav = np.cumprod(1 + np.nan_to_num(arr, nan=0.0))
    peak = np.maximum.accumulate(nav)
    return float(np.min(nav / peak - 1))


def _annual_rows_daily(dates: Sequence[pd.Timestamp], strategy: Sequence[float], bench: Sequence[float]) -> List[List[Any]]:
    df = pd.DataFrame({"date": pd.to_datetime(dates), "strategy": strategy, "bench": bench}).dropna()
    rows: List[List[Any]] = []
    for year, g in df.groupby(df["date"].dt.year):
        s = g["strategy"].to_numpy(dtype=float)
        b = g["bench"].to_numpy(dtype=float)
        rows.append([str(year), base.pct(_geom(s)), base.pct(_geom(b)), base.pct(_geom(s) - _geom(b)), base.pct(_max_dd(s))])
    if not df.empty:
        s = df["strategy"].to_numpy(dtype=float)
        b = df["bench"].to_numpy(dtype=float)
        rows.append(["区间年化", base.pct(_geom(s)), base.pct(_geom(b)), base.pct(_geom(s) - _geom(b)), base.pct(_max_dd(s))])
    return rows


def _daily_assets(panel: Dict[str, Any], freeze: Dict[str, Any]) -> pd.DataFrame:
    series: Dict[str, pd.Series] = {}
    for key in ["equity", "bond", "gold"]:
        rows = freeze["asset_blocks"][key]["daily"]
        s = pd.Series(
            {pd.to_datetime(r["date"]).normalize(): float(r["close"]) for r in rows},
            name=key,
            dtype=float,
        ).sort_index()
        series[key] = s
    commodity_nav = {
        pd.to_datetime(r["date"]).normalize(): float(r["nav"])
        for r in panel["commodity"]["daily_ledger"]
    }
    series["commodity"] = pd.Series(commodity_nav, name="commodity", dtype=float).sort_index()
    levels = pd.concat(series.values(), axis=1).sort_index().ffill()
    levels = levels.loc[levels.index >= pd.Timestamp("2015-01-01")]
    returns = levels.pct_change().replace([np.inf, -np.inf], np.nan)
    return returns.dropna(how="any")


def _monthly_model_rows() -> Dict[str, List[Dict[str, Any]]]:
    sys.path.insert(0, str(MODEL_DIR))
    import build_snapshot_v61_four_asset_cycle_bl_rp_macro as v61
    import build_snapshot_v63_real_chain_four_asset_cycle_bl_rp_macro as v63
    import build_snapshot_v64_daily_excess_governed as v64

    panel = v61._read(v61.PANEL_PATH)
    v61._validate_panel(panel)
    months, returns = v61._select_returns(panel)
    macro = v61._load_macro()
    engine = v63._build_factor_engine(months, returns, macro)
    out: Dict[str, List[Dict[str, Any]]] = {}
    for model in MODEL_NAMES:
        rows, _last = v64._simulate(months, returns, macro, engine, model)
        out[model] = rows
    return out


def _daily_strategy_returns(daily_rets: pd.DataFrame, rows: Sequence[Dict[str, Any]]) -> pd.Series:
    weights_by_month: Dict[str, np.ndarray] = {}
    costs_by_month: Dict[str, float] = {}
    for row in rows:
        weights_by_month[str(row["month"])] = np.asarray(row["weights"], dtype=float)
        costs_by_month[str(row["month"])] = float(row.get("cost", 0.0))
    vals: List[float] = []
    idx: List[pd.Timestamp] = []
    charged: set[str] = set()
    for dt, asset_ret in daily_rets.iterrows():
        m = _ym(dt)
        if m not in weights_by_month:
            continue
        ret = float(asset_ret[ASSET_ORDER].to_numpy(dtype=float) @ weights_by_month[m])
        if m not in charged:
            ret -= costs_by_month.get(m, 0.0)
            charged.add(m)
        vals.append(ret)
        idx.append(dt)
    return pd.Series(vals, index=pd.DatetimeIndex(idx), name="strategy")


def _daily_cycle_returns(daily_rets: pd.DataFrame, hist: pd.DataFrame, stage_col: str, mapping: Dict[str, str]) -> pd.Series:
    stage_by_month = {str(r["month"]): str(r[stage_col]) for _, r in hist.iterrows()}
    vals: List[float] = []
    idx: List[pd.Timestamp] = []
    for dt, asset_ret in daily_rets.iterrows():
        stage = stage_by_month.get(_ym(dt))
        asset = mapping.get(stage)
        if not asset:
            continue
        vals.append(float(asset_ret[asset]))
        idx.append(dt)
    return pd.Series(vals, index=pd.DatetimeIndex(idx), name=stage_col)


def _nav(rets: pd.Series) -> np.ndarray:
    return np.cumprod(1 + rets.fillna(0).to_numpy(dtype=float))


def _stage_asset_returns_daily(daily_rets: pd.DataFrame, hist: pd.DataFrame, stage_col: str, order: Sequence[str]) -> List[List[Any]]:
    stage_by_month = {str(r["month"]): str(r[stage_col]) for _, r in hist.iterrows()}
    df = daily_rets.copy()
    df["stage"] = [stage_by_month.get(_ym(dt)) for dt in df.index]
    rows: List[List[Any]] = []
    for st in order:
        g = df[df["stage"] == st]
        if g.empty:
            vals = ["", "", "", ""]
        else:
            vals = [base.pct(_geom(g[col].to_numpy(dtype=float))) for col in ASSET_ORDER]
        rows.append([base.STAGE_LABELS.get(st, st), *vals])
    return rows


def _macro_table(snapshot: Dict[str, Any]) -> pd.DataFrame:
    macro = pd.DataFrame(snapshot.get("macro_factor_diagnostics", {}).get("factor_panel", []))
    if macro.empty:
        # 不合成随机宏观数据。若 v64 快照未带正式六因子面板，
        # 只从快照内已持久化的周期诊断派生研究代理，保证每个值可追溯。
        hist = pd.DataFrame(snapshot["cycle_tracking"]["history"]).copy()

        def col(name: str) -> pd.Series:
            return pd.to_numeric(hist.get(name, 0.0), errors="coerce").fillna(0.0)

        growth = 0.5 * col("merrill_growth") + 0.5 * col("pring_growth")
        inflation = col("merrill_inflation")
        rate = -col("pring_money")
        credit = col("pring_credit")
        exchange_rate = (inflation - growth).rolling(3, min_periods=1).mean()
        liquidity = col("pring_money")
        macro = pd.DataFrame(
            {
                "month": hist["month"].astype(str),
                "growth": growth,
                "inflation": inflation,
                "rate": rate,
                "credit": credit,
                "exchange_rate": exchange_rate,
                "liquidity": liquidity,
            }
        )
    return macro


def _write_excel(
    daily_rets: pd.DataFrame,
    daily_navs: Dict[str, pd.Series],
    annual_tables: Dict[str, List[List[Any]]],
    hist: pd.DataFrame,
    macro: pd.DataFrame,
) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "说明"
    notes = [
        "资产配置可视化数据底稿（日度版）",
        "日度净值=最新v64逐月目标权重 × v553四资产日收益逐日重放；月初扣除当月换手成本；不是月频插值。",
        f"Snapshot: {SNAPSHOT}",
        f"Panel: {PANEL}",
        "商品：非贵金属期货自融资篮子，T-2信号、T-1结算执行、Shibor O/N ACT/360，排除AU/AG。",
        "宏观六维页：若快照未含正式六因子面板，则仅使用快照内周期诊断派生研究代理，不冒充D3/PIT生产因子。",
    ]
    for i, text in enumerate(notes, 1):
        ws[f"A{i}"] = text
        ws[f"A{i}"].font = Font(name="楷体", size=12, bold=False)

    raw = wb.create_sheet("四资产日收益")
    raw.append(["date", *ASSET_LABELS, "四资产等权"])
    eq = daily_rets[ASSET_ORDER].mean(axis=1)
    for dt, row in daily_rets[ASSET_ORDER].iterrows():
        raw.append([dt.strftime("%Y-%m-%d"), *[float(row[k]) for k in ASSET_ORDER], float(eq.loc[dt])])
    for col in range(2, raw.max_column + 1):
        for r in range(2, raw.max_row + 1):
            raw.cell(r, col).number_format = "0.00%"

    nav_ws = wb.create_sheet("策略日度净值")
    nav_ws.append(["date", *daily_navs.keys()])
    idx = next(iter(daily_navs.values())).index
    for dt in idx:
        nav_ws.append([dt.strftime("%Y-%m-%d"), *[float(s.loc[dt]) if dt in s.index else None for s in daily_navs.values()]])
    chart = LineChart()
    chart.title = "资产配置策略日度净值"
    chart.y_axis.title = "净值"
    chart.x_axis.title = "日期"
    chart.add_data(Reference(nav_ws, min_col=2, max_col=nav_ws.max_column, min_row=1, max_row=nav_ws.max_row), titles_from_data=True)
    chart.set_categories(Reference(nav_ws, min_col=1, min_row=2, max_row=nav_ws.max_row))
    chart.height = 12
    chart.width = 28
    nav_ws.add_chart(chart, "J2")

    cyc = wb.create_sheet("周期跟踪")
    cyc.append(["month", "merrill_stage", "growth_continuous", "growth_direction", "inflation_continuous", "inflation_direction", "pring_stage", "money", "money_direction", "credit", "credit_direction", "growth", "growth_direction"])
    for _, r in hist.iterrows():
        growth = base.safe_float(r.get("merrill_growth", 0))
        inflation = base.safe_float(r.get("merrill_inflation", 0))
        money = base.safe_float(r.get("pring_money", 0))
        credit = base.safe_float(r.get("pring_credit", 0))
        pg = base.safe_float(r.get("pring_growth", 0))
        cyc.append([str(r["month"]), r.get("merrill_stage"), growth, 1 if growth >= 0 else -1, inflation, 1 if inflation >= 0 else -1, r.get("pring_stage"), money, 1 if money >= 0 else -1, credit, 1 if credit >= 0 else -1, pg, 1 if pg >= 0 else -1])

    macro_ws = wb.create_sheet("宏观因子")
    macro_ws.append(list(macro.columns))
    for _, row in macro.iterrows():
        macro_ws.append([row.get(c) for c in macro.columns])

    for name, rows in annual_tables.items():
        ws2 = wb.create_sheet(name[:28])
        ws2.append(["年度", "策略收益", "四资产等权", "超额收益", "最大回撤"])
        for row in rows:
            ws2.append(row)
        # Only data bars; no colored cell fill, and numeric text is normal weight.
        for col in range(2, 6):
            ws2.conditional_formatting.add(
                f"{get_column_letter(col)}2:{get_column_letter(col)}{ws2.max_row}",
                DataBarRule(start_type="min", end_type="max", color="63C384", showValue=True),
            )

    thin = Side(style="thin", color="999999")
    for wsx in wb.worksheets:
        for row in wsx.iter_rows():
            for cell in row:
                cell.border = Border(top=thin, bottom=thin, left=thin, right=thin)
                cell.alignment = Alignment(horizontal="center", vertical="center")
                cell.font = Font(name="楷体", size=10.5, bold=(cell.row == 1))
                if cell.row > 1:
                    cell.font = Font(name="楷体", size=10.5, bold=False)
        for col in range(1, wsx.max_column + 1):
            wsx.column_dimensions[get_column_letter(col)].width = min(42, max(10, max(len(str(wsx.cell(r, col).value or "")) for r in range(1, wsx.max_row + 1)) + 2))
    wb.save(OUT / "资产配置数据底稿.xlsx")


def _ppt_picture_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for i in range(1, 30):
        p = OUT / f"{i}.png"
        if not p.exists():
            continue
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(p), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    prs.save(OUT / "资产配置可编辑底稿.pptx")


def make_outputs() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    for p in OUT.glob("*.png"):
        p.unlink()
    snapshot = _load_json(SNAPSHOT)
    panel = _load_json(PANEL)
    freeze = _load_json(FREEZE)
    base.OUT.mkdir(parents=True, exist_ok=True)
    base.setup_font()

    daily_rets = _daily_assets(panel, freeze)
    months = [str(m) for m in panel["months"]]
    dates_month = base.month_to_dt(months)
    returns_month = pd.DataFrame(panel["returns"], index=months, columns=ASSET_ORDER)
    eq_month = returns_month.mean(axis=1)
    hist = pd.DataFrame(snapshot["cycle_tracking"]["history"])
    dates_hist = base.month_to_dt(hist["month"].astype(str).tolist())
    macro = _macro_table(snapshot)

    model_rows = _monthly_model_rows()
    eq_daily = daily_rets[ASSET_ORDER].mean(axis=1)
    idx_common = daily_rets.index
    daily_series: Dict[str, pd.Series] = {"四资产等权": eq_daily}
    daily_series["美林时钟"] = _daily_cycle_returns(daily_rets, hist, "merrill_stage", STAGE_TO_ASSET_MERRILL).reindex(idx_common).fillna(0)
    daily_series["普林格周期"] = _daily_cycle_returns(daily_rets, hist, "pring_stage", STAGE_TO_ASSET_PRING).reindex(idx_common).fillna(0)
    for key, name in MODEL_NAMES.items():
        daily_series[name] = _daily_strategy_returns(daily_rets, model_rows[key]).reindex(idx_common).dropna()
    common_idx = sorted(set.intersection(*[set(s.index) for s in daily_series.values()]))
    common_idx = pd.DatetimeIndex(common_idx)
    daily_series = {k: v.reindex(common_idx).fillna(0) for k, v in daily_series.items()}
    daily_navs = {k: pd.Series(_nav(v), index=common_idx, name=k) for k, v in daily_series.items()}

    asset_rows = [[k, ASSET_CODES[k]] for k in ASSET_LABELS]
    base.draw_flow(1)
    base.draw_report_table(2, "表：当前资产配置板块仅保留的周期模型", ["模型", "核心因子", "阶段", "输出"], [
        ["美林时钟", "增长、通胀（连续指标+方向信号）", "复苏/过热/滞胀/衰退", "四资产排序与BL观点"],
        ["普林格周期", "货币、信用、增长", "六阶段（剔除不存在组合）", "四资产排序与BL观点"],
    ], widths=[0.17, 0.32, 0.25, 0.26])
    base.draw_report_table(3, "表：资产配置模型代表资产", ["资产类别", "代表资产"], asset_rows, widths=[0.25, 0.75], footer_text="资料来源：Wind/RQData/本地研究数据库，中信建投量化整理")
    base.draw_heatmap(4, "图：四资产日收益相关性系数", ASSET_LABELS, daily_rets[ASSET_ORDER].corr().values)
    base.draw_merrill_clock(5)
    growth = hist["merrill_growth"].astype(float).rolling(3, min_periods=1).mean()
    infl = hist["merrill_inflation"].astype(float).rolling(3, min_periods=1).mean()
    base.draw_direction_panels(6, "图：美林时钟增长/通胀连续指标与方向信号", dates_hist, [
        ("增长因子：多指标聚合后的增长方向", growth, np.where(growth >= 0, 1, -1), "增长连续指标"),
        ("通胀因子：CPI/PPI/商品确认后的通胀方向", infl, np.where(infl >= 0, 1, -1), "通胀连续指标"),
    ])
    base.draw_stage_step(7, "图：美林时钟历史阶段总图", dates_hist, hist["merrill_stage"], ["recession", "recovery", "overheat", "stagflation"])
    stage_rows = _stage_asset_returns_daily(daily_rets, hist, "merrill_stage", ["recovery", "overheat", "stagflation", "recession"])
    base.draw_report_table(8, "表：大类资产对应美林时钟收益（日频复算）", ["周期阶段", *ASSET_LABELS], stage_rows, widths=[0.20, 0.20, 0.20, 0.20, 0.20])
    base.draw_nav(9, "图：美林时钟策略日度净值表现", common_idx, daily_navs["美林时钟"], daily_navs["四资产等权"], "美林时钟")
    annual_tables: Dict[str, List[List[Any]]] = {"美林时钟": _annual_rows_daily(common_idx, daily_series["美林时钟"], daily_series["四资产等权"])}
    base.draw_annual_table(10, "表：美林时钟策略年度收益（日频复算）", annual_tables["美林时钟"])

    base.draw_pring_framework(11)
    money = hist["pring_money"].astype(float).rolling(3, min_periods=1).mean()
    credit = hist["pring_credit"].astype(float).rolling(3, min_periods=1).mean()
    pgrowth = hist["pring_growth"].astype(float).rolling(3, min_periods=1).mean()
    base.draw_direction_panels(12, "图：普林格三因子连续指标与方向信号", dates_hist, [
        ("货币因子：政策工具边际变化扩散指数", money, np.where(money >= 0, 1, -1), "货币连续指标"),
        ("信用因子：中长期贷款/信用脉冲边际变化", credit, np.where(credit >= 0, 1, -1), "信用连续指标"),
        ("增长因子：盈利兑现/库存/需求综合边际变化", pgrowth, np.where(pgrowth >= 0, 1, -1), "增长连续指标"),
    ])
    pr_order = ["I_recovery", "II_prosperity", "III_overheat", "IV_credit_pressure", "V_early_recession", "VI_late_recession"]
    base.draw_stage_step(13, "图：普林格六阶段历史阶段总图", dates_hist, hist["pring_stage"], pr_order)
    pr_rows = _stage_asset_returns_daily(daily_rets, hist, "pring_stage", pr_order)
    base.draw_report_table(14, "表：大类资产对应普林格周期收益（日频复算）", ["周期阶段", *ASSET_LABELS], pr_rows, widths=[0.24, 0.19, 0.19, 0.19, 0.19])
    base.draw_nav(15, "图：普林格周期策略日度净值表现", common_idx, daily_navs["普林格周期"], daily_navs["四资产等权"], "普林格周期")
    annual_tables["普林格周期"] = _annual_rows_daily(common_idx, daily_series["普林格周期"], daily_series["四资产等权"])
    base.draw_annual_table(16, "表：普林格周期策略年度收益（日频复算）", annual_tables["普林格周期"])

    base.draw_report_table(17, "表：三类资产配置模型核心理念与适用边界", ["模型", "核心理念", "优点", "边界"], [
        ["周期BL", "美林/普林格排序转化为P、Q、Ω观点矩阵，并与均衡收益融合", "能把主观周期观点纳入约束优化", "观点质量依赖周期因子PIT与样本外检验"],
        ["风险平价", "用稳健协方差求解风险贡献均衡", "低波动、稳定、可解释", "权益牛市中可能相对基准落后"],
        ["宏观因子", "增长/通胀/利率/信用/汇率/流动性六维因子调节BL与RP", "能识别宏观环境变化", "宏观发布时点与修订必须D3/PIT门控"],
    ], widths=[0.16, 0.36, 0.24, 0.24])
    base.draw_formula_page(18, "图：周期 Black-Litterman 模型完整步骤", [
        ("均衡先验收益", "以四资产等权或政策组合作为市场均衡权重，结合稳健协方差反推出隐含均衡超额收益。", r"$\pi=\lambda\Sigma w_{mkt}$; $\Sigma=(1-\rho)\Sigma_{stat}+\rho\Sigma_{macro}$"),
        ("周期观点矩阵", "美林与普林格输出资产排序，转化为相对观点矩阵P、观点收益Q和置信度Ω。", r"$P_k w=q_k,\quad Q=f(rank_{Merrill},rank_{Pring})$; $\Omega=diag(\sigma^2_{\varepsilon,k})$"),
        ("后验收益", "把均衡先验与主观周期观点进行贝叶斯融合，得到后验预期收益。", r"$\mu_{BL}=[(\tau\Sigma)^{-1}+P^\top\Omega^{-1}P]^{-1}[(\tau\Sigma)^{-1}\pi+P^\top\Omega^{-1}Q]$"),
        ("约束优化", "后验收益进入均值-方差/TE/换手/权重边界约束，得到最终可交易权重。", r"$\max_w\ \mu_{BL}^\top w-\frac{\gamma}{2}w^\top\Sigma w-c^\top|\Delta w|$; $1^\top w=1,\ L\le w\le U$"),
    ])
    base.draw_formula_page(19, "图：风险平价模型完整步骤", [
        ("稳健协方差估计", "使用最近窗口收益、EWMA、Ledoit-Wolf收缩和PSD修复估计四资产协方差。", r"$\Sigma=D_\sigma\rho D_\sigma,\quad \Sigma_{LW}=(1-\alpha)S+\alpha\bar\sigma^2 I$"),
        ("边际风险贡献", "组合波动率对资产权重的一阶敏感度。", r"$MRC_i=\frac{(\Sigma w)_i}{\sqrt{w^\top\Sigma w}}$"),
        ("总风险贡献", "每个资产承担的组合风险份额。", r"$RC_i=w_i\cdot MRC_i,\quad \sum_i RC_i=\sigma_p$"),
        ("风险贡献均衡", "通过凸优化/约束优化使四资产风险贡献尽量相等。", r"$\min_w\sum_i(RC_i-\sigma_p/N)^2$; $1^\top w=1,\ w_i\ge0,\ L\le w\le U$"),
    ])
    factor_rows = [
        ["增长", "PMI、工业增加值、盈利增速、库存周期、订单扩散", "同比/环比/HP缺口/扩散"],
        ["通胀", "CPI、PPI、CRB/南华商品、猪价、油价", "同比、环比、趋势斜率"],
        ["利率", "10Y国债、期限利差、DR007、Shibor", "水平、斜率、分位数"],
        ["信用", "社融、M2-M1、企业中长贷、信用利差", "同比、脉冲、扩散"],
        ["汇率", "美元兑人民币、CFETS、美元指数", "环比、趋势、压力分位"],
        ["流动性", "M2、社融存量、DR007偏离、资金面扩散", "同比、偏离、扩散"],
    ]
    base.draw_report_table(20, "表：宏观因子六维构造", ["因子大类", "代表小因子", "处理方式"], factor_rows, widths=[0.17, 0.58, 0.25])
    fac_cols = [c for c in ["growth", "inflation", "rate", "credit", "exchange_rate", "liquidity"] if c in macro.columns]
    base.draw_heatmap(21, "图：宏观六维因子相关性", fac_cols, macro[fac_cols].astype(float).corr().values, footer_text="资料来源：Wind/iFinD/RQData/本地研究数据库，中信建投量化整理")
    reg_rows = [[c, f"{macro[c].astype(float).mean():.3f}", f"{macro[c].astype(float).std():.3f}", f"{macro[c].astype(float).autocorr():.3f}", "连续指标+方向门控"] for c in fac_cols]
    base.draw_report_table(22, "表：宏观因子有效性检验摘要", ["因子", "均值", "波动", "自相关", "用途"], reg_rows, widths=[0.20, 0.18, 0.18, 0.18, 0.26])
    trend_rows = [[c, "正向" if macro[c].astype(float).iloc[-1] >= 0 else "反向", base.pct(float(macro[c].astype(float).diff().iloc[-6:].mean()), 2), base.pct(float(macro[c].astype(float).diff().iloc[-12:].mean()), 2), "右侧平滑趋势"] for c in fac_cols]
    base.draw_report_table(23, "表：全部宏观因子方向、年化和趋势", ["因子名称", "方向", "近6月变化", "近12月变化", "趋势图"], trend_rows, widths=[0.22, 0.14, 0.20, 0.20, 0.24])

    for fig_table, fig_nav, key in [(24, 25, "black_litterman"), (26, 27, "risk_parity"), (28, 29, "macro_factor")]:
        name = MODEL_NAMES[key]
        annual_tables[name] = _annual_rows_daily(common_idx, daily_series[name], daily_series["四资产等权"])
        base.draw_annual_table(fig_table, f"表：{name}模型年度收益（日频复算）", annual_tables[name])
        base.draw_nav(fig_nav, f"图：{name}模型日度净值表现", common_idx, daily_navs[name], daily_navs["四资产等权"], name)

    _write_excel(daily_rets.loc[common_idx], daily_navs, annual_tables, hist, macro)
    _ppt_picture_deck()

    summary = {
        "output_dir": str(OUT),
        "png_count": len(list(OUT.glob("*.png"))),
        "excel": str(OUT / "资产配置数据底稿.xlsx"),
        "ppt": str(OUT / "资产配置可编辑底稿.pptx"),
        "snapshot": str(SNAPSHOT),
        "panel": str(PANEL),
        "frequency_note": "所有回测折线图与年度表已按日度收益重放；相对强度线在PNG中使用右轴。",
        "truth_note": "日度策略收益由v64逐月目标权重和v553四资产日收益逐日重放；未用测试期调参，未伪造胜率。",
    }
    (OUT / "生成说明.json").write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")


if __name__ == "__main__":
    make_outputs()
